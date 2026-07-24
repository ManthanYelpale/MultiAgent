import logging
import os
import uuid
from datetime import datetime
from typing import Any

import numpy as np
import pandas as pd

logger = logging.getLogger(__name__)

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except ImportError:
    plt = None

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
except ImportError:
    colors = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

from app.core.config import settings
from app.services.ml_service import detect_anomalies, generate_ai_insights, run_forecast

REPORTS_DIR = os.path.join(settings.UPLOAD_DIR, "reports")


def generate_chart_image(df: pd.DataFrame, target_column: str, date_column: str | None = None) -> str | None:
    if plt is None or target_column not in df.columns:
        return None

    os.makedirs(REPORTS_DIR, exist_ok=True)
    chart_filename = f"chart_{uuid.uuid4().hex[:8]}.png"
    chart_path = os.path.join(REPORTS_DIR, chart_filename)

    plt.figure(figsize=(6.5, 3.2), dpi=150)
    series = pd.to_numeric(df[target_column], errors="coerce").dropna()

    if date_column and date_column in df.columns:
        x_vals = df.loc[series.index, date_column].astype(str)
        plt.plot(x_vals, series, marker="o", color="#4F46E5", linewidth=2)
        plt.xticks(rotation=30, ha="right", fontsize=8)
    else:
        plt.plot(range(1, len(series) + 1), series, marker="o", color="#4F46E5", linewidth=2)
        plt.xlabel("Index / Period", fontsize=8)

    plt.title(f"Metric Trend: {target_column}", fontsize=10, fontweight="bold")
    plt.ylabel(target_column, fontsize=8)
    plt.grid(True, linestyle="--", alpha=0.5)
    plt.tight_layout()
    plt.savefig(chart_path)
    plt.close()

    return chart_path


def build_pdf_report(
    pdf_path: str,
    original_filename: str,
    summary_stats: dict[str, Any],
    insights_text: str,
    forecast_data: dict[str, Any] | None = None,
    chart_path: str | None = None,
) -> str:
    if SimpleDocTemplate is None:
        raise RuntimeError("reportlab is not installed.")

    os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#1E1B4B"),
        spaceAfter=6,
    )
    subtitle_style = ParagraphStyle(
        "SubTitle",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        textColor=colors.HexColor("#6B7280"),
        spaceAfter=15,
    )
    heading_style = ParagraphStyle(
        "Heading",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=16,
        textColor=colors.HexColor("#312E81"),
        spaceBefore=12,
        spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#1F2937"),
        spaceAfter=10,
    )

    story = []
    # Title & Subtitle
    story.append(Paragraph("AI Executive Weekly Performance Report", title_style))
    date_str = datetime.now().strftime("%B %d, %Y")
    story.append(Paragraph(f"Dataset: <b>{original_filename}</b> | Date Generated: {date_str}", subtitle_style))
    story.append(Spacer(1, 10))

    # Executive Insights Paragraph
    story.append(Paragraph("1. Executive AI Insights Summary", heading_style))
    story.append(Paragraph(insights_text, body_style))
    story.append(Spacer(1, 10))

    # Matplotlib Visual Chart
    if chart_path and os.path.exists(chart_path):
        story.append(Paragraph("2. Performance Visualizations", heading_style))
        story.append(Image(chart_path, width=6*inch, height=2.85*inch))
        story.append(Spacer(1, 10))

    # Statistical & Anomaly Summary Table
    story.append(Paragraph("3. Dataset Metrics & Anomaly Audit", heading_style))
    table_data = [
        ["Metric Category", "Value / Count"],
        ["Total Rows Processed", str(summary_stats.get("total_rows", "N/A"))],
        ["Flagged Anomalies", str(summary_stats.get("anomaly_count", 0))],
        ["Data Columns", ", ".join(summary_stats.get("columns", [])[:5])],
    ]
    t = Table(table_data, colWidths=[2.5*inch, 4*inch])
    t.setStyle(
        TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4F46E5")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F9FAFB")),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E5E7EB")),
        ])
    )
    story.append(t)
    story.append(Spacer(1, 10))

    # Metric Forecast Table if available
    if forecast_data and "forecast" in forecast_data:
        story.append(Paragraph("4. Metric Projections (Forecast)", heading_style))
        fc_table_data = [["Period", f"Projected {forecast_data.get('target_column', 'Value')}"]]
        for item in forecast_data["forecast"]:
            fc_table_data.append([str(item.get("period_label")), str(item.get("forecast_value"))])

        fc_table = Table(fc_table_data, colWidths=[3*inch, 3.5*inch])
        fc_table.setStyle(
            TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#065F46")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D1D5DB")),
            ])
        )
        story.append(fc_table)

    doc.build(story)
    return pdf_path


def build_pptx_report(
    pptx_path: str,
    original_filename: str,
    summary_stats: dict[str, Any],
    insights_text: str,
    forecast_data: dict[str, Any] | None = None,
    chart_path: str | None = None,
) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed.")

    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    prs = Presentation()

    # Slide 1: Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)

    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "AI Executive Weekly Business Report"
    p1.font.bold = True
    p1.font.size = Pt(28)
    p1.font.color.rgb = RGBColor(30, 27, 75)

    p2 = tf1.add_paragraph()
    p2.text = f"Dataset: {original_filename} | Generated: {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(107, 114, 128)

    # Slide 2: Executive Insights & Summary Metrics
    slide2 = prs.slides.add_slide(blank_slide_layout)
    tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1))
    tf2 = tb2.text_frame
    p_s2 = tf2.paragraphs[0]
    p_s2.text = "Executive AI Insights & Summary"
    p_s2.font.bold = True
    p_s2.font.size = Pt(22)
    p_s2.font.color.rgb = RGBColor(49, 46, 129)

    tb2_body = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5))
    tf2_body = tb2_body.text_frame
    tf2_body.word_wrap = True
    p_body = tf2_body.paragraphs[0]
    p_body.text = insights_text
    p_body.font.size = Pt(12)

    p_stats = tf2_body.add_paragraph()
    p_stats.text = (
        f"\nMetrics Summary:\n"
        f"• Total Rows: {summary_stats.get('total_rows', 'N/A')}\n"
        f"• Flagged Anomalies: {summary_stats.get('anomaly_count', 0)}\n"
        f"• Key Columns: {', '.join(summary_stats.get('columns', [])[:5])}"
    )
    p_stats.font.size = Pt(12)
    p_stats.font.color.rgb = RGBColor(79, 70, 229)

    # Slide 3: Chart & Metric Forecast
    slide3 = prs.slides.add_slide(blank_slide_layout)
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1))
    tf3 = tb3.text_frame
    p_s3 = tf3.paragraphs[0]
    p_s3.text = "Metric Visualizations & Forecasting"
    p_s3.font.bold = True
    p_s3.font.size = Pt(22)
    p_s3.font.color.rgb = RGBColor(49, 46, 129)

    if chart_path and os.path.exists(chart_path):
        slide3.shapes.add_picture(chart_path, Inches(0.8), Inches(1.5), width=Inches(5.0))

    if forecast_data and "forecast" in forecast_data:
        tb_fc = slide3.shapes.add_textbox(Inches(6.0), Inches(1.5), Inches(3.5), Inches(4))
        tf_fc = tb_fc.text_frame
        p_fc_title = tf_fc.paragraphs[0]
        p_fc_title.text = f"Forecast ({forecast_data.get('target_column')}):"
        p_fc_title.font.bold = True
        p_fc_title.font.size = Pt(13)

        for item in forecast_data["forecast"]:
            p_item = tf_fc.add_paragraph()
            p_item.text = f"• {item.get('period_label')}: {item.get('forecast_value')}"
            p_item.font.size = Pt(11)

    prs.save(pptx_path)
    return pptx_path


def generate_weekly_report_pipeline(
    db_file: Any,
    df: pd.DataFrame,
    target_column: str | None = None,
    date_column: str | None = None,
    export_format: str = "both",
) -> dict[str, Any]:
    os.makedirs(REPORTS_DIR, exist_ok=True)
    report_id = uuid.uuid4().hex[:8]

    # Select numerical target column if not specified
    if not target_column:
        num_cols = df.select_dtypes(include=[np.number]).columns
        target_column = num_cols[0] if len(num_cols) > 0 else df.columns[0]

    # Run anomaly detection
    anomalies_res = detect_anomalies(df)

    # Run forecasting
    forecast_res = None
    try:
        forecast_res = run_forecast(df, target_column=target_column, date_column=date_column, horizon=5)
    except Exception as exc:
        logger.warning(f"Skipping forecast generation: {exc}")

    # Generate AI insights
    insights_text = generate_ai_insights(
        filename=db_file.original_filename,
        df=df,
        anomalies_info=anomalies_res,
        forecast_info=forecast_res,
    )

    # Render Matplotlib Chart
    chart_path = generate_chart_image(df, target_column=target_column, date_column=date_column)

    summary_stats = {
        "total_rows": len(df),
        "anomaly_count": anomalies_res.get("anomaly_count", 0),
        "columns": list(df.columns),
    }

    pdf_filename = None
    pptx_filename = None

    if export_format in ["pdf", "both"]:
        pdf_filename = f"report_{db_file.id}_{report_id}.pdf"
        pdf_path = os.path.join(REPORTS_DIR, pdf_filename)
        try:
            build_pdf_report(
                pdf_path=pdf_path,
                original_filename=db_file.original_filename,
                summary_stats=summary_stats,
                insights_text=insights_text,
                forecast_data=forecast_res,
                chart_path=chart_path,
            )
        except Exception as exc:
            logger.error(f"Error building PDF: {exc}")
            pdf_filename = None

    if export_format in ["pptx", "both"]:
        pptx_filename = f"presentation_{db_file.id}_{report_id}.pptx"
        pptx_path = os.path.join(REPORTS_DIR, pptx_filename)
        try:
            build_pptx_report(
                pptx_path=pptx_path,
                original_filename=db_file.original_filename,
                summary_stats=summary_stats,
                insights_text=insights_text,
                forecast_data=forecast_res,
                chart_path=chart_path,
            )
        except Exception as exc:
            logger.error(f"Error building PPTX: {exc}")
            pptx_filename = None

    return {
        "file_id": db_file.id,
        "original_filename": db_file.original_filename,
        "insights_summary": insights_text,
        "pdf_filename": pdf_filename,
        "pptx_filename": pptx_filename,
    }
