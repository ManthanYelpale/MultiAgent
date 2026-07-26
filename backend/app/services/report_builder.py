import logging
import os
import uuid
from datetime import datetime
from typing import Any, List, Dict

logger = logging.getLogger(__name__)

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Image as RLImage, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle, PageBreak
except ImportError:
    # Every name from the block must be bound, otherwise the `if SimpleDocTemplate is
    # None` guard below raises NameError instead of the intended RuntimeError.
    colors = letter = ParagraphStyle = getSampleStyleSheet = inch = None
    RLImage = Paragraph = SimpleDocTemplate = Spacer = Table = TableStyle = PageBreak = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

from app.core.config import settings
from app.services.chart_renderer import render_chart_to_image

REPORTS_DIR = os.path.join(settings.UPLOAD_DIR, "reports")

class ReportBuilder:
    def __init__(self, original_filename: str, sections: List[str]):
        self.original_filename = original_filename
        self.sections = sections
        self.date_str = datetime.now().strftime("%B %d, %Y")
        os.makedirs(REPORTS_DIR, exist_ok=True)
        self.report_id = uuid.uuid4().hex[:8]

    def build_pdf(
        self,
        pdf_filename: str,
        kpis: List[Dict[str, Any]] = None,
        insights_text: str = None,
        charts: List[Dict[str, Any]] = None,
        forecast_data: dict = None
    ) -> str:
        if SimpleDocTemplate is None:
            raise RuntimeError("reportlab is not installed.")

        pdf_path = os.path.join(REPORTS_DIR, pdf_filename)
        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=letter,
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36,
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "DocTitle",
            parent=styles["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=24,
            leading=28,
            textColor=colors.HexColor("#1E1B4B"),
            spaceAfter=12,
            alignment=1, # Center
        )
        subtitle_style = ParagraphStyle(
            "SubTitle",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=12,
            textColor=colors.HexColor("#6B7280"),
            spaceAfter=30,
            alignment=1, # Center
        )
        heading_style = ParagraphStyle(
            "Heading",
            parent=styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=16,
            leading=20,
            textColor=colors.HexColor("#312E81"),
            spaceBefore=20,
            spaceAfter=10,
        )
        body_style = ParagraphStyle(
            "Body",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=11,
            leading=16,
            textColor=colors.HexColor("#1F2937"),
            spaceAfter=10,
        )

        story = []

        # 1. Cover Page
        story.append(Spacer(1, 150))
        story.append(Paragraph("Data Analytics Executive Report", title_style))
        story.append(Paragraph(f"Dataset: <b>{self.original_filename}</b><br/>Date: {self.date_str}", subtitle_style))
        story.append(PageBreak())

        # 2. KPIs (if selected)
        if "kpis" in self.sections and kpis:
            story.append(Paragraph("Key Performance Indicators", heading_style))
            # Format KPI data
            table_data = [["Metric", "Value"]]
            for kpi in kpis:
                table_data.append([kpi.get("title", ""), str(kpi.get("value", ""))])
            
            t = Table(table_data, colWidths=[3*inch, 3*inch])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4F46E5")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 11),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F9FAFB")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E5E7EB")),
            ]))
            story.append(t)
            story.append(Spacer(1, 20))
            
            # KPI AI Insights below the table
            for kpi in kpis:
                if kpi.get("insight"):
                    story.append(Paragraph(f"<b>{kpi.get('title')}</b>: ✨ {kpi.get('insight')}", body_style))

            story.append(PageBreak())

        # 3. AI Insights (if selected)
        if "insights" in self.sections and insights_text:
            story.append(Paragraph("Executive AI Insights", heading_style))
            # Split by newlines if necessary to render correctly
            for para in insights_text.split("\n"):
                if para.strip():
                    story.append(Paragraph(para.strip(), body_style))
            story.append(PageBreak())

        # 4. Charts (if selected)
        if "charts" in self.sections and charts:
            story.append(Paragraph("Dashboard Visualizations", heading_style))
            for chart in charts:
                chart_type = chart.get("chart_type")
                title = f"{chart.get('agg_function', '').capitalize()} of {chart.get('y_column', '')} by {chart.get('x_column', '')}"
                if chart_type == "pie":
                    title = f"Distribution of {chart.get('x_column', '')}"
                
                story.append(Paragraph(title, ParagraphStyle("ChartTitle", parent=heading_style, fontSize=14)))
                
                # Render chart to PNG
                img_path = render_chart_to_image(chart_type, chart.get("data", []), title)
                if img_path:
                    story.append(RLImage(img_path, width=6*inch, height=3.2*inch))
                    story.append(Spacer(1, 15))
                
                # If the chart has any specific insights, add them here
                if chart.get("insight"):
                    story.append(Paragraph(f"<b>Insight:</b> {chart.get('insight')}", body_style))
                
                story.append(Spacer(1, 25))

        # 5. Forecast (if selected)
        if "forecast" in self.sections and forecast_data and "forecast" in forecast_data:
            story.append(PageBreak())
            story.append(Paragraph("Metric Projections (Forecast)", heading_style))
            
            target_col = forecast_data.get('target_column', 'Value')
            fc_table_data = [["Period", f"Projected {target_col}"]]
            for item in forecast_data["forecast"]:
                fc_table_data.append([str(item.get("period_label")), str(item.get("forecast_value"))])

            fc_table = Table(fc_table_data, colWidths=[3*inch, 3*inch])
            fc_table.setStyle(
                TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#065F46")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 11),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D1D5DB")),
                ])
            )
            story.append(fc_table)

        doc.build(story)
        return pdf_path

    def build_pptx(
        self,
        pptx_filename: str,
        kpis: List[Dict[str, Any]] = None,
        insights_text: str = None,
        charts: List[Dict[str, Any]] = None,
        forecast_data: dict = None
    ) -> str:
        if Presentation is None:
            raise RuntimeError("python-pptx is not installed.")

        pptx_path = os.path.join(REPORTS_DIR, pptx_filename)
        prs = Presentation()

        # Helper to add titles
        def add_slide_title(slide, text):
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(49, 46, 129)

        # 1. Title Slide
        blank_slide_layout = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(blank_slide_layout)

        tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(2))
        tf1 = tb1.text_frame
        p1 = tf1.paragraphs[0]
        p1.text = "Data Analytics Executive Report"
        p1.font.bold = True
        p1.font.size = Pt(32)
        p1.font.color.rgb = RGBColor(30, 27, 75)

        p2 = tf1.add_paragraph()
        p2.text = f"Dataset: {self.original_filename} | Generated: {self.date_str}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(107, 114, 128)

        # 2. KPIs
        if "kpis" in self.sections and kpis:
            slide_kpi = prs.slides.add_slide(blank_slide_layout)
            add_slide_title(slide_kpi, "Key Performance Indicators")
            
            top = Inches(1.5)
            for kpi in kpis:
                tb = slide_kpi.shapes.add_textbox(Inches(0.5), top, Inches(4), Inches(0.8))
                tf = tb.text_frame
                tf.text = f"{kpi.get('title')}: {kpi.get('value')}"
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(18)
                
                if kpi.get("insight"):
                    p_insight = tf.add_paragraph()
                    p_insight.text = f"✨ {kpi.get('insight')}"
                    p_insight.font.size = Pt(12)
                    p_insight.font.color.rgb = RGBColor(107, 114, 128)
                
                top += Inches(1.2)

        # 3. AI Insights
        if "insights" in self.sections and insights_text:
            slide_ins = prs.slides.add_slide(blank_slide_layout)
            add_slide_title(slide_ins, "Executive AI Insights")
            
            tb_body = slide_ins.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
            tf_body = tb_body.text_frame
            tf_body.word_wrap = True
            p_body = tf_body.paragraphs[0]
            p_body.text = insights_text
            p_body.font.size = Pt(14)

        # 4. Charts
        if "charts" in self.sections and charts:
            for chart in charts:
                slide_chart = prs.slides.add_slide(blank_slide_layout)
                
                chart_type = chart.get("chart_type")
                title = f"{chart.get('agg_function', '').capitalize()} of {chart.get('y_column', '')} by {chart.get('x_column', '')}"
                if chart_type == "pie":
                    title = f"Distribution of {chart.get('x_column', '')}"
                
                add_slide_title(slide_chart, title)
                
                # Render chart to PNG
                img_path = render_chart_to_image(chart_type, chart.get("data", []), title)
                if img_path:
                    slide_chart.shapes.add_picture(img_path, Inches(1.0), Inches(1.5), width=Inches(8.0))
                
                if chart.get("insight"):
                    tb_insight = slide_chart.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(1.0))
                    tf_insight = tb_insight.text_frame
                    tf_insight.word_wrap = True
                    p_ins = tf_insight.paragraphs[0]
                    p_ins.text = f"Insight: {chart.get('insight')}"
                    p_ins.font.size = Pt(12)
                    p_ins.font.color.rgb = RGBColor(79, 70, 229)

        # 5. Forecast
        if "forecast" in self.sections and forecast_data and "forecast" in forecast_data:
            slide_fc = prs.slides.add_slide(blank_slide_layout)
            target_col = forecast_data.get('target_column', 'Value')
            add_slide_title(slide_fc, f"Metric Projections ({target_col})")
            
            tb_fc = slide_fc.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
            tf_fc = tb_fc.text_frame
            
            for item in forecast_data["forecast"]:
                p_item = tf_fc.add_paragraph()
                p_item.text = f"• {item.get('period_label')}: {item.get('forecast_value')}"
                p_item.font.size = Pt(16)

        prs.save(pptx_path)
        return pptx_path
